"""Build Gate 2 presentation for Richa Sang — Apex Distribution Ltd."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1B, 0x2A, 0x4A)   # headings / backgrounds
TEAL    = RGBColor(0x00, 0x7A, 0x87)   # accent
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF0, 0xF4, 0xF8)   # slide background
MID     = RGBColor(0x4A, 0x5D, 0x75)   # body text
GREEN   = RGBColor(0x2E, 0x86, 0x48)
AMBER   = RGBColor(0xD4, 0x7E, 0x00)
RED     = RGBColor(0xC0, 0x39, 0x2B)

W  = Inches(13.33)   # widescreen width
H  = Inches(7.5)     # widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # truly blank layout


# ── Helper functions ─────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    fill  = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = Pt(line_width_pt)
    else:
        line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                word_wrap=True, italic=False):
    txb  = slide.shapes.add_textbox(x, y, w, h)
    tf   = txb.text_frame
    tf.word_wrap = word_wrap
    p    = tf.paragraphs[0]
    p.alignment = align
    run  = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color or MID
    return txb


def header_bar(slide, title, subtitle=None):
    """Dark navy bar across the top."""
    add_rect(slide, 0, 0, W, Inches(1.25), fill_rgb=NAVY)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.12), Inches(11), Inches(0.65),
                font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.4), Inches(0.75), Inches(11), Inches(0.42),
                    font_size=14, color=TEAL, italic=True)


def slide_bg(slide):
    add_rect(slide, 0, 0, W, H, fill_rgb=LIGHT)


def bullet_block(slide, items, x, y, w, h,
                 font_size=15, color=None, bullet="▸ "):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}{item}"
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color or MID


def card(slide, x, y, w, h, title, body_items,
         header_rgb=None, title_size=15, body_size=13):
    header_rgb = header_rgb or TEAL
    add_rect(slide, x, y, w, Inches(0.42), fill_rgb=header_rgb)
    add_textbox(slide, title,
                x + Inches(0.1), y + Inches(0.04), w - Inches(0.2), Inches(0.36),
                font_size=title_size, bold=True, color=WHITE)
    add_rect(slide, x, y + Inches(0.42), w, h - Inches(0.42),
             fill_rgb=WHITE, line_rgb=RGBColor(0xCC, 0xD6, 0xE0), line_width_pt=1)
    bullet_block(slide, body_items,
                 x + Inches(0.12), y + Inches(0.48),
                 w - Inches(0.24), h - Inches(0.58),
                 font_size=body_size)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill_rgb=NAVY)
# accent strip
add_rect(sl, 0, Inches(4.5), W, Inches(0.06), fill_rgb=TEAL)

add_textbox(sl, "Gate 2 Submission",
            Inches(1), Inches(1.2), Inches(11), Inches(0.8),
            font_size=18, bold=False, color=TEAL, align=PP_ALIGN.CENTER)
add_textbox(sl, "Agentic Transformation of\nApex Distribution Ltd",
            Inches(1), Inches(1.9), Inches(11), Inches(1.8),
            font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(sl, "Customer Operations — 35 people · 730 cases/day · 4 work streams",
            Inches(1), Inches(3.6), Inches(11), Inches(0.55),
            font_size=17, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)
add_textbox(sl, "Richa Sang  ·  2026-05-06",
            Inches(1), Inches(5.9), Inches(11), Inches(0.45),
            font_size=14, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)

add_textbox(sl, "7 Deliverables  ·  Honest AI Assessment  ·  Phase 1 Recommendation",
            Inches(1), Inches(5.2), Inches(11), Inches(0.45),
            font_size=13, color=TEAL, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Situation
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "The Situation", "What Apex is dealing with — and why previous AI attempts failed")

# 4 work stream cards
streams = [
    ("ETA Inquiries",        ["400/day · 4 min avg", "55% of all cases", "Simple data lookup", "✅ High automation potential"],   GREEN),
    ("Delivery Exceptions",  ["180/day · 12 min avg", "Damage, refusals, missed windows", "Incomplete SOP (Section 4.3 blank)", "⚠️  Low automation ceiling"],  AMBER),
    ("Dispatch Adjustments", ["90/day · 18 min avg", "Mid-route re-routing", "Citrix barrier — no AI integration", "🚫 Not accessible now"],          RED),
    ("Billing Disputes",     ["60/day · 28 min avg", "Mostly Aurum system friction", "T-1 batch only — no real-time", "⚠️  Partial automation only"],   AMBER),
]
cw = Inches(2.9)
gap = Inches(0.2)
start_x = Inches(0.4)
for i, (title, items, col) in enumerate(streams):
    card(sl, start_x + i*(cw+gap), Inches(1.45), cw, Inches(3.8),
         title, items, header_rgb=col, body_size=13)

# Prior failures box
add_rect(sl, Inches(0.4), Inches(5.45), Inches(12.5), Inches(1.7),
         fill_rgb=WHITE, line_rgb=RED, line_width_pt=2)
add_textbox(sl, "Sarah Whitmore's prior AI attempts (both failed):",
            Inches(0.6), Inches(5.55), Inches(12), Inches(0.35),
            font_size=13, bold=True, color=RED)
add_textbox(sl,
    "▸ 2023 chatbot — customers rejected it; tried to cover too many work streams at once\n"
    "▸ 2024 automated billing tool — broke every time Aurum schema changed (no version control on quarterly exports)",
    Inches(0.6), Inches(5.9), Inches(12), Inches(1.0),
    font_size=12, color=MID)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — 7 Deliverables Overview
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "7 Deliverables — What Was Built", "ATX Gate 2 methodology — each deliverable drives the next")

deliverables = [
    ("D1", "Cognitive Load Map",           "How all 4 work streams actually happen — JTBDs, 6 cognitive dimensions, zone maps"),
    ("D2", "Delegation Matrix",            "Scored suitability for AI-led vs. human-led per work stream"),
    ("D3", "Volume × Value Analysis",      "Where AI investment delivers highest return for lowest risk — with financial model"),
    ("D4", "Agent Purpose Document",       "Developer-ready spec for the ETA agent — inputs, outputs, escalation logic, KPIs"),
    ("D5", "System / Data Inventory",      "What systems exist, what's accessible, what's blocked, what's risky"),
    ("D6", "Discovery Questions",          "11 questions whose answers would change what gets built — surfaces real unknowns"),
    ("D7", "CLAUDE.md",                    "Working project file for AI coding agent — enforces design decisions in build phase"),
]

row_h = Inches(0.72)
for i, (code, title, desc) in enumerate(deliverables):
    y = Inches(1.45) + i * row_h
    add_rect(sl, Inches(0.4), y, Inches(0.8), row_h - Inches(0.06),
             fill_rgb=NAVY)
    add_textbox(sl, code, Inches(0.4), y + Inches(0.12), Inches(0.8), Inches(0.48),
                font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(1.25), y, Inches(11.65), row_h - Inches(0.06),
             fill_rgb=WHITE, line_rgb=RGBColor(0xCC, 0xD6, 0xE0), line_width_pt=1)
    add_textbox(sl, title, Inches(1.4), y + Inches(0.1), Inches(2.8), Inches(0.52),
                font_size=14, bold=True, color=NAVY)
    add_textbox(sl, desc, Inches(4.3), y + Inches(0.1), Inches(8.4), Inches(0.52),
                font_size=13, color=MID)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — D1: Cognitive Load Map
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "D1 — Cognitive Load Map", "All 4 work streams decomposed: JTBDs · 6 cognitive dimensions · Zone maps · Breakpoints")

cols = [
    ("ETA Inquiries\n(400/day · 4 min)",
     ["Zone A: ~65% (fully agentic)", "Zone B: ~30% (agent-led)", "Working memory: LOW", "Pattern recognition: HIGH", "Judgment: LOW", "Time pressure: LOW", "→ Primary automation target"],
     GREEN),
    ("Delivery Exceptions\n(180/day · 12 min)",
     ["Zone A: ~0%", "Zone B: ~15% (context only)", "Working memory: VERY HIGH", "Judgment: VERY HIGH", "Emotional labour: HIGH", "Incomplete SOP (Sec 4.3 blank)", "→ Human-led; AI assists"],
     AMBER),
    ("Dispatch Adjustments\n(90/day · 18 min)",
     ["Zone A: ~0% (Citrix barrier)", "Zone B: 10–15% (context surfacing)", "Working memory: VERY HIGH", "Time pressure: VERY HIGH", "Pattern recognition: HIGH (uncoded)", "Multi-vehicle cascade risk", "→ Human-only until re-platform"],
     RED),
    ("Billing Disputes\n(60/day · 28 min)",
     ["Zone A: 0% (Aurum ceiling)", "Zone B: 30–35% (triage + drafting)", "28 min = 5–8 min decisions", "+ 20+ min Aurum friction", "Audit bypass risk (Sandra)", "→ Human-led; AI assists in triage"],
     TEAL),
]
cw = Inches(2.95)
gap = Inches(0.18)
for i, (title, items, col) in enumerate(cols):
    card(sl, Inches(0.35) + i*(cw+gap), Inches(1.4), cw, Inches(5.7),
         title, items, header_rgb=col, title_size=13, body_size=12)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — D2: Delegation Suitability Matrix
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "D2 — Delegation Suitability Matrix", "6-dimension scored assessment · Automation readiness score · Archetype assignment")

# Table header
cols_w = [Inches(2.8), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.6)]
headers = ["Work Stream", "Data Quality", "Judgment Req.", "Error Risk", "SOP Maturity", "System Access", "Volume", "Verdict"]
col_colors = [NAVY, NAVY, NAVY, NAVY, NAVY, NAVY, NAVY, NAVY]

y0 = Inches(1.45)
x0 = Inches(0.3)
row_h = Inches(0.55)

for j, (hdr, cw) in enumerate(zip(headers, cols_w)):
    x = x0 + sum(cols_w[:j])
    add_rect(sl, x, y0, cw, row_h, fill_rgb=NAVY)
    add_textbox(sl, hdr, x+Inches(0.05), y0+Inches(0.08), cw-Inches(0.1), row_h-Inches(0.1),
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("ETA Inquiries",        "High",   "Low",     "Low",     "Complete", "REST API ✅",  "400/day", ("AI leads", GREEN)),
    ("Delivery Exceptions",  "Medium", "V.High",  "High",    "Incomplete","Accessible ✅","180/day", ("Human leads", RED)),
    ("Dispatch Adjustments", "Low",    "V.High",  "V.High",  "Partial",  "Citrix 🚫",   "90/day",  ("Human-only", RED)),
    ("Billing Disputes",     "Low",    "Medium",  "High",    "Partial",  "Batch-only ⚠","60/day",  ("Human leads", AMBER)),
]
row_fills = [RGBColor(0xE8,0xF5,0xEB), RGBColor(0xFF,0xF3,0xE0),
             RGBColor(0xFF,0xEB,0xEE), RGBColor(0xFF,0xF3,0xE0)]

for i, (ws, dq, jr, er, sop, sa, vol, verdict) in enumerate(rows):
    y = y0 + (i+1)*row_h
    bg = row_fills[i]
    cells = [ws, dq, jr, er, sop, sa, vol, verdict[0]]
    for j, (cell, cw) in enumerate(zip(cells, cols_w)):
        x = x0 + sum(cols_w[:j])
        fc = bg if j < 7 else verdict[1]
        tc = MID if j < 7 else WHITE
        add_rect(sl, x, y, cw, row_h,
                 fill_rgb=fc, line_rgb=RGBColor(0xCC,0xD6,0xE0), line_width_pt=0.5)
        add_textbox(sl, cell, x+Inches(0.05), y+Inches(0.1), cw-Inches(0.1), row_h-Inches(0.12),
                    font_size=12, bold=(j==0 or j==7), color=tc, align=PP_ALIGN.CENTER)

# Key insight
add_rect(sl, Inches(0.3), Inches(4.1), Inches(12.7), Inches(0.9),
         fill_rgb=WHITE, line_rgb=TEAL, line_width_pt=2)
add_textbox(sl,
    "Key insight: 2 of 4 work streams score below 3/5 for automation readiness. "
    "This is named openly — the biggest failure pattern in AI projects is assuming everything can be automated.",
    Inches(0.5), Inches(4.18), Inches(12.3), Inches(0.72),
    font_size=13, color=NAVY, italic=True)

# 6 dimensions note
bullet_block(sl,
    ["6 dimensions scored 1–5: Data Quality, Judgment Requirement (inverse), Error Risk (inverse), SOP Maturity, System Accessibility, Volume",
     "3 dimensions are inverse-scored — higher automation need = lower score on those axes"],
    Inches(0.3), Inches(5.15), Inches(12.7), Inches(1.0),
    font_size=12, color=MID)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — D3: Volume × Value + ROI
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "D3 — Volume × Value Analysis + ROI", "Where to start and what it's actually worth — honest numbers")

# Quadrant visual (text-based)
add_rect(sl, Inches(0.35), Inches(1.4), Inches(5.5), Inches(5.7),
         fill_rgb=WHITE, line_rgb=RGBColor(0xCC,0xD6,0xE0), line_width_pt=1)
add_textbox(sl, "Volume × Value Quadrant",
            Inches(0.5), Inches(1.5), Inches(5.2), Inches(0.4),
            font_size=13, bold=True, color=NAVY)

# axes labels
add_textbox(sl, "HIGH VALUE ↑",
            Inches(0.4), Inches(1.9), Inches(1.2), Inches(0.35), font_size=10, color=MID, italic=True)
add_textbox(sl, "LOW VALUE",
            Inches(0.4), Inches(6.5), Inches(1.2), Inches(0.35), font_size=10, color=MID, italic=True)
add_textbox(sl, "LOW VOLUME ←————————→ HIGH VOLUME",
            Inches(0.5), Inches(6.9), Inches(5.2), Inches(0.3), font_size=10, color=MID, italic=True)

quads = [
    (Inches(3.2), Inches(2.1),  "ETA Inquiries",        "Vol: HIGH / Value: HIGH", GREEN,  "★ PRIMARY"),
    (Inches(3.0), Inches(3.8),  "Billing Disputes",     "Vol: LOW / Value: MED",   TEAL,   "Phase 2"),
    (Inches(1.0), Inches(3.6),  "Delivery Exceptions",  "Vol: MED / Value: MED",   AMBER,  "Phase 3"),
    (Inches(1.0), Inches(5.5),  "Dispatch Adjustments", "Vol: LOW / Value: LOW",   RED,    "Phase 4"),
]
for qx, qy, name, sub, col, tag in quads:
    add_rect(sl, qx, qy, Inches(2.3), Inches(0.75), fill_rgb=col)
    add_textbox(sl, f"{tag}: {name}", qx+Inches(0.06), qy+Inches(0.04), Inches(2.2), Inches(0.35),
                font_size=11, bold=True, color=WHITE)
    add_textbox(sl, sub, qx+Inches(0.06), qy+Inches(0.38), Inches(2.2), Inches(0.3),
                font_size=10, color=WHITE)

# ROI panel
add_rect(sl, Inches(6.1), Inches(1.4), Inches(6.9), Inches(5.7),
         fill_rgb=WHITE, line_rgb=TEAL, line_width_pt=2)
add_textbox(sl, "Phase 1 ROI — ETA Inquiries Agent",
            Inches(6.3), Inches(1.5), Inches(6.5), Inches(0.4),
            font_size=14, bold=True, color=NAVY)

roi_lines = [
    ("Recoverable minutes/day",   "~1,120 min  (70% of 1,600)"),
    ("Person-hours/day",           "~18.7 hrs"),
    ("FTE-equivalent capacity",    "~2.5 FTE"),
    ("Birmingham agent fully-loaded", "~£33,000/yr\n(salary + NI + pension + overhead)"),
    ("Released capacity value",    "~£82,500/yr"),
]
y_roi = Inches(2.0)
for label, value in roi_lines:
    add_textbox(sl, label, Inches(6.3), y_roi, Inches(3.3), Inches(0.5),
                font_size=12, color=MID)
    add_textbox(sl, value, Inches(9.7), y_roi, Inches(3.1), Inches(0.5),
                font_size=12, bold=True, color=NAVY)
    y_roi += Inches(0.52)

# divider
add_rect(sl, Inches(6.3), y_roi + Inches(0.05), Inches(6.5), Inches(0.03), fill_rgb=TEAL)
y_roi += Inches(0.18)

add_textbox(sl, "vs CEO benchmark of £1.2M",
            Inches(6.3), y_roi, Inches(6.5), Inches(0.35),
            font_size=12, bold=True, color=RED)
y_roi += Inches(0.38)
add_textbox(sl,
    "The £1.2M competitor saving reflects broader scope + headcount reduction + "
    "more modern systems. Phase 1 (ETA only) = £82,500. Phases 1–3 combined "
    "= ~£200,000–300,000/yr. The gap must be named with the CEO — "
    "not hidden. See D6 Q9.",
    Inches(6.3), y_roi, Inches(6.5), Inches(1.4),
    font_size=11, color=MID)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — D4: Agent Purpose Document
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "D4 — Agent Purpose Document", "Developer-ready spec for the ETA Agent — one job, done well")

# Flow diagram (text boxes + arrows)
steps = [
    ("Inbound message\n(SMS / email / web portal)", TEAL),
    ("Order resolver\n(Salesforce CRM REST API)", NAVY),
    ("GPS interpreter\n(Driver App — ⚠ API not yet confirmed)", AMBER),
    ("Autonomy evaluator\n(rules engine: GPS age + flags)", NAVY),
    ("Response composer\nor Escalation router", GREEN),
    ("CRM case logger\n(every interaction — mandatory)", NAVY),
]
box_w = Inches(1.95)
box_h = Inches(0.9)
gap_x = Inches(0.22)
total_w = len(steps)*box_w + (len(steps)-1)*gap_x
start_x = (W - total_w) / 2
y_flow = Inches(1.6)

for i, (label, col) in enumerate(steps):
    x = start_x + i*(box_w+gap_x)
    add_rect(sl, x, y_flow, box_w, box_h, fill_rgb=col)
    add_textbox(sl, label, x+Inches(0.06), y_flow+Inches(0.1), box_w-Inches(0.12), box_h-Inches(0.15),
                font_size=11, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        add_textbox(sl, "→", x+box_w, y_flow+Inches(0.3), gap_x, Inches(0.35),
                    font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Critical constraints
add_textbox(sl, "Critical constraints (hard-coded, non-negotiable):",
            Inches(0.4), Inches(2.75), Inches(12.5), Inches(0.38),
            font_size=13, bold=True, color=NAVY)

constraints = [
    ("GPS < 30 min", "Autonomous — reply with best-estimate window", GREEN),
    ("GPS 30–60 min", "Supervised — widened window + disclosure; flagged", AMBER),
    ("GPS > 60 min", "Escalate — NO ETA given. Dispatcher notified.", RED),
    ("Driver App offline", "Escalate — never respond. Immediate escalation.", RED),
    ("Exception flag set", "Escalate — no ETA regardless of GPS freshness", RED),
    ("Strategic account\n(ACCT_MGR set)", "Escalate to named queue — validated in APEX_CUSTOMER_MASTER", RED),
]
cw2 = Inches(2.05)
gap2 = Inches(0.12)
for i, (trigger, action, col) in enumerate(constraints):
    x = Inches(0.35) + i*(cw2+gap2)
    card(sl, x, Inches(3.2), cw2, Inches(3.85),
         trigger, [action], header_rgb=col, title_size=12, body_size=11)

# Scope boundary note
add_rect(sl, Inches(0.35), Inches(7.1), Inches(12.65), Inches(0.28),
         fill_rgb=NAVY)
add_textbox(sl,
    "Scope: ETA inquiries ONLY. Not a chatbot. Cannot handle exceptions, billing, complaints, route changes, or requests to speak to a person.",
    Inches(0.5), Inches(7.12), Inches(12.4), Inches(0.25),
    font_size=11, bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — D5: System / Data Inventory
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "D5 — System / Data Inventory", "What's accessible · What's blocked · What's risky")

systems = [
    ("✅ Salesforce CRM",
     ["REST API confirmed", "Primary order lookup", "Case logging target", "Rate limits — confirm before batch"],
     GREEN),
    ("⚠️  Driver GPS App",
     ["In-house iOS/Android app", "GPS API NOT CONFIRMED", "Gate-blocking question", "Build blocked until confirmed"],
     AMBER),
    ("🚫 Dispatch Console\n(Citrix)",
     ["No clean API surface", "Citrix = RPA-class only (brittle)", "Major infra project to integrate", "Phase 4 — not now"],
     RED),
    ("⚠️  Aurum Billing\n(2008 Oracle)",
     ["Batch CSV export only", "T-1 minimum data lag", "48h to modify invoices", "Quarterly schema changes — no notice"],
     RED),
]
cw3 = Inches(2.9)
gap3 = Inches(0.2)
for i, (title, items, col) in enumerate(systems):
    card(sl, Inches(0.35) + i*(cw3+gap3), Inches(1.4), cw3, Inches(4.3),
         title, items, header_rgb=col, title_size=13, body_size=12)

# Audit risk callout
add_rect(sl, Inches(0.35), Inches(5.9), Inches(12.65), Inches(1.4),
         fill_rgb=RGBColor(0xFF,0xEB,0xEE), line_rgb=RED, line_width_pt=2)
add_textbox(sl, "Billing Audit Risk — must fix BEFORE any AI touches billing",
            Inches(0.55), Inches(5.98), Inches(12.2), Inches(0.38),
            font_size=13, bold=True, color=RED)
add_textbox(sl,
    "Artefact 2 (Billing case study): Sandra applied a £170 goodwill credit via a manual workaround — no APPROVER_ID, "
    "no AUDIT_REF, no audit trail entry. The formal approval process exists on paper but is being bypassed. "
    "Building AI on top of this broken audit trail compounds the control gap. Fix governance first.",
    Inches(0.55), Inches(6.38), Inches(12.2), Inches(0.85),
    font_size=12, color=MID)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — D6: Discovery Questions
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "D6 — Discovery Questions", "11 questions that would change what gets built — surfaces genuine unknowns")

questions = [
    ("Q1 — GATE-BLOCKING",
     "Does the Driver App backend expose a REST API for GPS telemetry?\nIf NO: precision ETA feature cannot be built.",
     RED),
    ("Q2 — Design-changing",
     "What specifically made the 2024 chatbot fail — the channel (customers didn't want a bot) "
     "or the answers (wrong information)? Changes how the new system is designed and presented.",
     AMBER),
    ("Q3 — Governance",
     "How often do team members bypass the formal credit approval process? Is Sandra's £170 credit an "
     "isolated incident or a team-wide pattern? Determines whether Phase 2 is a software or governance project first.",
     AMBER),
    ("Q4 — Calibration",
     "How stale is GPS data in practice? The Artefact showed a 26-minute gap. If typical, "
     "the agent will escalate more cases than planned — changing KPI projections.",
     AMBER),
    ("Q5 — Scope",
     "What is the channel split for the 400 ETA inquiries/day — SMS vs email vs web portal? "
     "Determines how many inbound channels the comms gateway must support.",
     TEAL),
]
row_h2 = Inches(1.05)
for i, (qnum, qtext, col) in enumerate(questions):
    y = Inches(1.45) + i * row_h2
    add_rect(sl, Inches(0.35), y, Inches(2.0), row_h2 - Inches(0.06), fill_rgb=col)
    add_textbox(sl, qnum, Inches(0.4), y+Inches(0.1), Inches(1.9), row_h2-Inches(0.18),
                font_size=12, bold=True, color=WHITE)
    add_rect(sl, Inches(2.4), y, Inches(10.6), row_h2 - Inches(0.06),
             fill_rgb=WHITE, line_rgb=RGBColor(0xCC,0xD6,0xE0), line_width_pt=1)
    add_textbox(sl, qtext, Inches(2.55), y+Inches(0.08), Inches(10.3), row_h2-Inches(0.18),
                font_size=12, color=MID)

add_textbox(sl,
    "Q6–Q9: Channel split · Billing active vs elapsed time · Sandra's authority structure · SOP gap  "
    "·  Q10: Contractual SLAs and breach tracking  ·  Q11: Salesforce ACCT_MGR / key-account flag "
    "(VALIDATED — build-ready in Escalation Trigger 9)",
    Inches(0.35), Inches(6.85), Inches(12.65), Inches(0.5),
    font_size=11, color=MID, italic=True)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — D7: CLAUDE.md
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "D7 — CLAUDE.md: Project File for Developers", "Enforces design decisions — prevents scope creep in the build phase")

left_items = [
    ("What it is",
     ["Working guide an AI coding agent picks up at build time",
      "Translates Deliverables 1–6 into concrete build instructions",
      "Enforces constraints that must not erode during development"]),
    ("Build order (recommended)",
     ["1. CRM logger — build and test first (everything feeds it)",
      "2. Order resolver — CRM lookup + fallback",
      "3. Autonomy evaluator — rules engine, unit-test before integration",
      "4. Response composer — channel-appropriate templates",
      "5. Escalation router — test every trigger before ship",
      "6. GPS interpreter — BLOCKED until Driver App API confirmed"]),
]

right_items = [
    ("Hard rules enforced in code",
     ["GPS > 60 min → NEVER give ETA, escalate",
      "Exception flag set → escalate regardless of GPS",
      "Delivery 'complete' with no scan event → escalate (missing parcel)",
      "Any billing question → route to billing team, do not process",
      "CRM log required for every interaction including escalations",
      "Never extrapolate GPS position from a stale ping"]),
    ("What it prevents",
     ["Scope creep: agent stays ETA-only, full stop",
      "Silent failures: schema-change detection required",
      "Audit gaps: billing-adjacent data must disclose T-1 lag",
      "Test shortcuts: 16 unit tests required before ship"]),
]

def left_right_cards(slide, left, right):
    for group, items_list in left:
        pass  # will build below

y_start = Inches(1.45)
col_w = Inches(6.1)
for col_i, group_list in enumerate([left_items, right_items]):
    x_base = Inches(0.35) + col_i * (col_w + Inches(0.4))
    y_cur = y_start
    for title, items in group_list:
        block_h = Inches(0.42) + len(items) * Inches(0.38)
        card(sl, x_base, y_cur, col_w, block_h, title, items,
             header_rgb=NAVY, title_size=13, body_size=12)
        y_cur += block_h + Inches(0.15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — What Happens Next
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "What Happens Next", "Sequenced programme — prerequisites respected, no wishful shortcuts")

phases = [
    ("NOW — Pre-build",
     ["Confirm Driver App GPS API (gate-blocker)",
      "Fix billing credit audit trail (governance before AI)",
      "Update Exception SOP — Section 4.3 is blank",
      "Schedule Salesforce schema review with Apex IT"],
     TEAL, "Immediate"),
    ("Phase 1\nETA Agent",
     ["Build CRM logger first",
      "Order resolver + autonomy evaluator",
      "Response composer + escalation router",
      "GPS interpreter (when API confirmed)",
      "Target: 70% of 400/day autonomous"],
     GREEN, "Buildable now"),
    ("Phase 2\nBilling Triage",
     ["Prerequisite: audit trail fixed",
      "Prerequisite: Aurum batch pipeline built",
      "Agent: triage + data retrieval + response draft",
      "Human approves all credits",
      "Hayes & Sons pattern detection"],
     AMBER, "After Phase 1 stable"),
    ("Phase 3\nException Context",
     ["Prerequisite: SOP updated",
      "Prerequisite: insurance protocol documented",
      "Agent: context assembly for dispatcher",
      "Human leads all decisions",
      "No autonomous exception resolution"],
     AMBER, "After Phase 2"),
    ("Phase 4\nDispatch Auto.",
     ["Prerequisite: Dispatch Console re-platform",
      "Prerequisite: API surface built",
      "Separate multi-year infrastructure project",
      "Not an AI project — a systems project",
      "Highest cascade risk of any stream"],
     RED, "Multi-year"),
]
cw4 = Inches(2.35)
gap4 = Inches(0.14)
for i, (title, items, col, tag) in enumerate(phases):
    x = Inches(0.3) + i*(cw4+gap4)
    add_rect(sl, x, Inches(1.4), cw4, Inches(0.38), fill_rgb=col)
    add_textbox(sl, tag, x+Inches(0.05), Inches(1.43), cw4-Inches(0.1), Inches(0.32),
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, x, Inches(1.78), cw4, Inches(0.5), fill_rgb=NAVY)
    add_textbox(sl, title, x+Inches(0.05), Inches(1.82), cw4-Inches(0.1), Inches(0.44),
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, x, Inches(2.28), cw4, Inches(4.8),
             fill_rgb=WHITE, line_rgb=RGBColor(0xCC,0xD6,0xE0), line_width_pt=1)
    bullet_block(sl, items, x+Inches(0.1), Inches(2.35), cw4-Inches(0.2), Inches(4.6),
                 font_size=11)

# ROI summary bar
add_rect(sl, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.3), fill_rgb=NAVY)
add_textbox(sl,
    "Phase 1: ~£82,500/yr released capacity  ·  "
    "Phases 1–3: ~£200,000–300,000/yr  ·  "
    "Full £1.2M requires modern system baseline + broader scope — name this gap with CEO",
    Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.27),
    font_size=11, bold=False, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Closing: The Honest Version
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill_rgb=NAVY)
add_rect(sl, 0, Inches(4.8), W, Inches(0.06), fill_rgb=TEAL)

add_textbox(sl, "The Honest Version",
            Inches(1), Inches(0.8), Inches(11), Inches(0.6),
            font_size=18, color=TEAL, align=PP_ALIGN.CENTER)
add_textbox(sl, "What this assessment actually says",
            Inches(1), Inches(1.35), Inches(11), Inches(0.7),
            font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

honests = [
    "Two of four work streams are not ready for AI automation — and this report says so plainly",
    "Phase 1 delivers ~£82,500/yr — not £1.2M — and the gap is explained, not hidden",
    "The gate-blocking question (Driver App API) is named before a single line of code is written",
    "The billing audit bypass must be fixed before AI is introduced — not after",
    "Narrow, testable, demonstrable — that is what a working AI programme looks like at this stage",
]
y_h = Inches(2.3)
for item in honests:
    add_textbox(sl, f"▸  {item}",
                Inches(1.5), y_h, Inches(10), Inches(0.48),
                font_size=14, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.LEFT)
    y_h += Inches(0.52)

add_textbox(sl, "Richa Sang  ·  Gate 2  ·  Apex Distribution Ltd  ·  2026-05-06",
            Inches(1), Inches(6.9), Inches(11), Inches(0.4),
            font_size=13, color=RGBColor(0x88,0x99,0xAA), align=PP_ALIGN.CENTER, italic=True)


# ── Save ────────────────────────────────────────────────────────────────────
out = r"c:\Users\RichaSang\gate-2-submission\Gate2-Richa-Sang-Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
